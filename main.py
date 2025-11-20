import asyncio
import logging
import json
import os
import requests
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import Command, StateFilter
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from aiogram.types import FSInputFile, InlineKeyboardMarkup, InlineKeyboardButton, WebAppInfo
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
from io import BytesIO

# Ma'lumotlar bazasi import
from models import (
    initialize_database, 
    get_user_by_telegram_id, 
    create_user, 
    get_user_statistics,
    User,
    Presentation as PresentationDB,
    ActivityLog
)

TELEGRAM_BOT_TOKEN = "8204596949:AAGWanbrVpFSGDh6BUI7YLJrmPIt4UcgN4Q"
GOOGLE_API_KEY = "AIzaSyDZRDi1t1bjtlOLmUvkWijtWEgI2m620HE"
WEBAPP_URL = "https://your-domain.com/registration.html"  

logging.basicConfig(level=logging.INFO)

bot = Bot(token=TELEGRAM_BOT_TOKEN)
dp = Dispatcher()

genai.configure(api_key=GOOGLE_API_KEY)
model = genai.GenerativeModel('gemini-2.0-flash')

# Ma'lumotlar bazasini ishga tushirish
initialize_database()


class PresentationForm(StatesGroup):
    waiting_for_language = State()
    waiting_for_slides = State()
    waiting_for_background = State()
    waiting_for_topic = State()


user_data = {}

BACKGROUND_COLORS = {
    "gradient_blue": [(25, 118, 210), (66, 165, 245)],
    "gradient_purple": [(123, 31, 162), (206, 147, 216)],
    "gradient_ocean": [(0, 150, 136), (0, 188, 212)],
    "gradient_sunset": [(255, 87, 34), (255, 193, 7)],
    "gradient_forest": [(27, 94, 32), (76, 175, 80)],
    "gradient_night": [(33, 33, 33), (97, 97, 97)],
    "gradient_pink": [(233, 30, 99), (244, 67, 54)],
    "gradient_gold": [(184, 134, 11), (255, 215, 0)]
}

BACKGROUND_NAMES = {
    "gradient_blue": "🔵 Osmon Ko'k",
    "gradient_purple": "💜 Binafshayi Qizg'in",
    "gradient_ocean": "🌊 Okean Suv",
    "gradient_sunset": "🌅 Quyosh Botish",
    "gradient_forest": "🌲 O'rmon Sabz",
    "gradient_night": "🌙 Tun Qorasi",
    "gradient_pink": "💗 Qizg'in Pushti",
    "gradient_gold": "✨ Oltin Nur"
}


def generate_presentation_content(topic, language, num_slides):
    """Gemini-dan kengaytirilgan prezentatsiya ma'lumotlarini JSON formatida oladi"""
    lang_name = "O'zbekcha" if language == "uz" else "English"
    
    prompt = f"""
    Sen professional prezentatsiya ustasisan. 
    Mavzu: "{topic}".
    Menga {num_slides} ta slayddan iborat JUDA KENGAYTIRILGAN prezentatsiya tayyorla.
    Har bir slayd: sarlavha, 6-8 ta BATAFSIL bullet point va rasm auksuli (rasm nomi).
    
    Javobni faqat va faqat toza JSON formatida qaytar:
    [
        {{
            "title": "Slayd sarlavhasi",
            "content": "• Batafsil nuqta 1\\n• Batafsil nuqta 2\\n• Batafsil nuqta 3\\n• Batafsil nuqta 4\\n• Batafsil nuqta 5\\n• Batafsil nuqta 6",
            "image_search": "rasm qidiruv so'zi (masalan: 'artificial intelligence')"
        }},
        ...
    ]
    
    Qoidalar:
    - {num_slides} ta slayd bo'lsin
    - Matn {lang_name} bo'lsin
    - Har bir slaydda 6-8 ta BATAFSIL bullet point bo'lsin
    - Har bir nuqta 1-2 ta jumla bo'lsin
    - image_search qismi 2-3 so'z bo'lsin
    - Faqat JSON jo'nating, boshqa hech narsa yo'q
    """
    
    try:
        response = model.generate_content(prompt)
        clean_text = response.text.replace("```json", "").replace("```", "").strip()
        return json.loads(clean_text)
    except Exception as e:
        logging.error(f"Xatolik Gemini: {e}")
        return None


def download_image(query):
    """Pixabay API orqali rasm yuklab oladi"""
    try:
        pixabay_api_key = "46608840-5b8e8d1b19e5a1e8c8e8d1b1"
        url = f"https://pixabay.com/api/?key={pixabay_api_key}&q={query}&image_type=photo&per_page=3&safesearch=true"
        
        response = requests.get(url, timeout=10)
        data = response.json()
        
        if data['hits']:
            image_url = data['hits'][0]['webformatURL']
            img_response = requests.get(image_url, timeout=10)
            return Image.open(BytesIO(img_response.content))
    except Exception as e:
        logging.error(f"Rasm yuklab olishda xatolik: {e}")
    
    return None


def create_gradient_image(width, height, color1, color2):
    """Gradient fon rasmi yaratadi"""
    img = Image.new('RGB', (width, height), color1)
    pixels = img.load()
    
    for y in range(height):
        ratio = y / height
        r = int(color1[0] * (1 - ratio) + color2[0] * ratio)
        g = int(color1[1] * (1 - ratio) + color2[1] * ratio)
        b = int(color1[2] * (1 - ratio) + color2[2] * ratio)
        
        for x in range(width):
            pixels[x, y] = (r, g, b)
    
    return img


def create_background_preview(bg_key):
    """Fon rangining preview rasmini yaratadi"""
    colors = BACKGROUND_COLORS[bg_key]
    img = create_gradient_image(400, 250, tuple(colors[0]), tuple(colors[1]))
    
    draw = ImageDraw.Draw(img)
    text = BACKGROUND_NAMES[bg_key]
    draw.text((20, 110), text, fill=(255, 255, 255))
    
    return img


def create_pptx(slides_data, user_id, bg_color_name):
    """JSON ma'lumotdan PPTX fayl yasaydi"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    colors = BACKGROUND_COLORS[bg_color_name]
    
    for idx, slide_info in enumerate(slides_data):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        gradient_img = create_gradient_image(1280, 960, tuple(colors[0]), tuple(colors[1]))
        img_path = f"temp_bg_{idx}.png"
        gradient_img.save(img_path)
        
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
        os.remove(img_path)
        
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = slide_info.get('title', 'Sarlavha')
        title_frame.word_wrap = True
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        try:
            image_query = slide_info.get('image_search', 'abstract')
            img = download_image(image_query)
            
            if img:
                img_path = f"temp_img_{idx}.png"
                img.save(img_path)
                slide.shapes.add_picture(img_path, Inches(6), Inches(1.1), width=Inches(3.8), height=Inches(2.8))
                os.remove(img_path)
        except Exception as e:
            logging.error(f"Rasm qo'shishda xatolik: {e}")
        
        content_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.1), Inches(5.5), Inches(6))
        content_frame = content_box.text_frame
        content_frame.text = slide_info.get('content', 'Matn yo\'q')
        content_frame.word_wrap = True
        
        for para in content_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.space_after = Pt(6)
            para.space_before = Pt(2)
    
    filename = f"presentation_{user_id}.pptx"
    prs.save(filename)
    return filename


@dp.message(Command("start"))
async def cmd_start(message: types.Message, state: FSMContext):
    user_id = message.from_user.id
    
    # Foydalanuvchi ro'yxatdan o'tganmi tekshirish
    user = get_user_by_telegram_id(user_id)
    
    if user is None:
        # Ro'yxatdan o'tmagan - WebApp ochish
        kb = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(
                text="📝 Ro'yxatdan o'tish",
                web_app=WebAppInfo(url=WEBAPP_URL)
            )]
        ])
        
        await message.answer(
            "👋 Salom! Botdan foydalanish uchun ro'yxatdan o'ting:\n\n"
            "🔐 Ma'lumotlaringiz xavfsiz saqlanadi\n"
            "⚡ Bir marta ro'yxatdan o'ting, doim foydalaning",
            reply_markup=kb
        )
    else:
        # Ro'yxatdan o'tgan - prezentatsiya yaratish
        user_data[user_id] = {}
        await state.set_state(PresentationForm.waiting_for_language)
        
        stats = get_user_statistics(user)
        
        kb = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="🇺🇿 O'zbekcha", callback_data="lang_uz")],
            [InlineKeyboardButton(text="🇬🇧 English", callback_data="lang_en")],
            [InlineKeyboardButton(text="📊 Statistikam", callback_data="show_stats")]
        ])
        
        await message.answer(
            f"Salom, {user.full_name}! 👋\n\n"
            f"📊 Siz {stats['total_presentations']} ta prezentatsiya yaratdingiz\n"
            f"📑 Jami {stats['total_slides']} ta slayd\n\n"
            "Tilni tanlang:",
            reply_markup=kb
        )


@dp.message(F.web_app_data)
async def handle_webapp_data(message: types.Message):
    """WebApp dan kelgan ma'lumotlarni qabul qilish"""
    try:
        data = json.loads(message.web_app_data.data)
        
        telegram_id = message.from_user.id
        full_name = data['full_name']
        gmail = data['gmail']
        phone = data['phone_number']
        age = data['age']
        
        # Foydalanuvchini bazaga qo'shish
        user = create_user(
            telegram_id=telegram_id,
            full_name=full_name,
            gmail=gmail,
            phone_number=phone,
            age=age
        )
        
        if user:
            await message.answer(
                f"✅ Xush kelibsiz, {full_name}!\n\n"
                f"📧 Gmail: {gmail}\n"
                f"📱 Telefon: {phone}\n"
                f"🎂 Yosh: {age}\n\n"
                f"Endi /start buyrug'ini bosing va prezentatsiya yarating!"
            )
        else:
            await message.answer(
                "❌ Xatolik yuz berdi. Bu email yoki Telegram ID allaqachon ro'yxatdan o'tgan.\n\n"
                "Agar muammo bo'lsa, @support ga murojaat qiling."
            )
            
    except Exception as e:
        logging.error(f"WebApp data xatolik: {e}")
        await message.answer("❌ Ma'lumotlarni qayta ishlashda xatolik yuz berdi.")


@dp.callback_query(F.data == "show_stats")
async def show_statistics(callback: types.CallbackQuery):
    """Foydalanuvchi statistikasini ko'rsatish"""
    user_id = callback.from_user.id
    user = get_user_by_telegram_id(user_id)
    
    if user:
        stats = get_user_statistics(user)
        
        recent_text = ""
        if stats['recent_presentations']:
            recent_text = "\n\n📋 Oxirgi prezentatsiyalar:\n"
            for pres in stats['recent_presentations']:
                recent_text += f"• {pres.topic} ({pres.slide_count} slayd)\n"
        
        await callback.message.answer(
            f"📊 Sizning statistikangiz:\n\n"
            f"👤 Ism: {user.full_name}\n"
            f"📧 Gmail: {user.gmail}\n"
            f"📱 Telefon: {user.phone_number}\n"
            f"🎂 Yosh: {user.age}\n\n"
            f"📈 Yaratilgan prezentatsiyalar: {stats['total_presentations']}\n"
            f"📑 Jami slaydlar: {stats['total_slides']}\n"
            f"📅 Ro'yxatdan o'tgan: {stats['registration_date']}\n"
            f"🕒 Oxirgi kirish: {stats['last_login']}"
            f"{recent_text}"
        )
    
    await callback.answer()


@dp.callback_query(StateFilter(PresentationForm.waiting_for_language))
async def process_language(callback: types.CallbackQuery, state: FSMContext):
    user_id = callback.from_user.id
    
    if callback.data == "lang_uz":
        user_data[user_id]['language'] = "uz"
        await callback.answer("O'zbekcha tanlandi ✓")
    else:
        user_data[user_id]['language'] = "en"
        await callback.answer("English selected ✓")
    
    await state.set_state(PresentationForm.waiting_for_slides)
    
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="5 slayd", callback_data="slides_5")],
        [InlineKeyboardButton(text="8 slayd", callback_data="slides_8")],
        [InlineKeyboardButton(text="10 slayd", callback_data="slides_10")],
        [InlineKeyboardButton(text="15 slayd", callback_data="slides_15")]
    ])
    
    await callback.message.edit_text("Slayd sonini tanlang:", reply_markup=kb)


@dp.callback_query(StateFilter(PresentationForm.waiting_for_slides))
async def process_slides(callback: types.CallbackQuery, state: FSMContext):
    user_id = callback.from_user.id
    
    slides_map = {"slides_5": 5, "slides_8": 8, "slides_10": 10, "slides_15": 15}
    user_data[user_id]['slides'] = slides_map[callback.data]
    await callback.answer(f"{user_data[user_id]['slides']} slayd tanlandi ✓")
    
    await state.set_state(PresentationForm.waiting_for_background)
    
    await callback.message.edit_text(
        f"✓ Slaydlar: {user_data[user_id]['slides']}\n\n🎨 Fon rangini tanlang:"
    )
    
    for bg_key in BACKGROUND_COLORS.keys():
        preview_img = create_background_preview(bg_key)
        img_path = f"preview_{bg_key}.png"
        preview_img.save(img_path)
        
        kb = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="✅ Tanlayman", callback_data=f"bg_{bg_key}")]
        ])
        
        await callback.message.chat.send_photo(
            photo=FSInputFile(img_path),
            caption=BACKGROUND_NAMES[bg_key],
            reply_markup=kb
        )
        os.remove(img_path)


@dp.callback_query(StateFilter(PresentationForm.waiting_for_background))
async def process_background(callback: types.CallbackQuery, state: FSMContext):
    user_id = callback.from_user.id
    
    bg_key = callback.data.replace("bg_", "")
    user_data[user_id]['background'] = bg_key
    await callback.answer("Fon tanlandi ✓")
    
    await state.set_state(PresentationForm.waiting_for_topic)
    
    await callback.message.answer(
        f"✓ Fon: {BACKGROUND_NAMES[bg_key]}\n\n📝 Endi mavzuni yozing:"
    )


@dp.message(StateFilter(PresentationForm.waiting_for_topic))
async def process_topic(message: types.Message, state: FSMContext):
    user_id = message.from_user.id
    user_topic = message.text
    
    # Foydalanuvchini olish
    user = get_user_by_telegram_id(user_id)
    
    lang = user_data[user_id]['language']
    num_slides = user_data[user_id]['slides']
    
    if lang == "uz":
        await message.answer(
            f"🔍 '{user_topic}' mavzusi bo'yicha ma'lumot yig'yapman...\n"
            f"📊 {num_slides} ta slayd, rasmlar bilan\n\n"
            f"⏳ Biroz kuting (30-40 soniya)..."
        )
    else:
        await message.answer(
            f"🔍 Gathering information about '{user_topic}'...\n"
            f"📊 {num_slides} slides with images\n\n"
            f"⏳ Please wait (30-40 seconds)..."
        )
    
    slides_data = generate_presentation_content(user_topic, lang, num_slides)
    
    if not slides_data:
        if lang == "uz":
            await message.answer("❌ Uzr, ma'lumotni shakllantirishda xatolik bo'ldi. /start bosib qayta urinib ko'ring.")
        else:
            await message.answer("❌ Sorry, error occurred. Try /start again.")
        await state.clear()
        return
    
    try:
        file_path = create_pptx(slides_data, user_id, user_data[user_id]['background'])
        ppt_file = FSInputFile(file_path)
        
        if lang == "uz":
            caption = f"✅ Mana, '{user_topic}' bo'yicha prezentatsiyangiz tayyor!\n\n📊 {num_slides} ta slayd + rasmlar"
        else:
            caption = f"✅ Your presentation about '{user_topic}' is ready!\n\n📊 {num_slides} slides + images"
        
        await message.answer_document(ppt_file, caption=caption)
        os.remove(file_path)
        
        # Statistikani yangilash
        if user:
            user.increment_presentations(num_slides)
            
            # Prezentatsiyani bazaga saqlash
            PresentationDB.create(
                user=user,
                topic=user_topic,
                language=lang,
                slide_count=num_slides,
                background_color=user_data[user_id]['background'],
                file_path=file_path
            )
            
            # Log yozish
            ActivityLog.create(
                user=user,
                action_type='create_presentation',
                description=f"Prezentatsiya yaratildi: {user_topic} ({num_slides} slayd)"
            )
        
        if lang == "uz":
            await message.answer("🎉 Yana prezentatsiya yasamoqchimisiz? /start bosing")
        else:
            await message.answer("🎉 Want to create another? Press /start")
        
    except Exception as e:
        logging.error(f"Xatolik PPTX: {e}")
        if lang == "uz":
            await message.answer("❌ Fayl yaratishda xatolik yuz berdi.")
        else:
            await message.answer("❌ File creation error.")
    
    await state.clear()


async def main():
    await dp.start_polling(bot)


if __name__ == "__main__":
    asyncio.run(main())